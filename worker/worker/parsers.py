

import io, pandas as pd
from pypdf import PdfReader
from docx import Document as Docx
from pptx import Presentation

def parse_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    texts = []
    for page in reader.pages:
        texts.append(page.extract_text() or "")
    return "\n".join(texts)

def parse_docx(data: bytes) -> str:
    doc = Docx(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs)

def parse_pptx(data: bytes) -> str:
    pres = Presentation(io.BytesIO(data))
    texts = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def parse_table(data: bytes, ext: str) -> str:
    if ext.lower() in (".xlsx", ".xls"):
        df = pd.read_excel(io.BytesIO(data))
    else:
        df = pd.read_csv(io.BytesIO(data))
    return df.to_csv(index=False)

def parse_txt(data: bytes) -> str:
    try:
        return data.decode("utf-8")
    except:
        return data.decode("latin-1", errors="ignore")
